from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
import pandas as pd
from datetime import date
import requests
from requests.auth import HTTPBasicAuth
from collections import defaultdict
import dateutil.parser
from pptx.dml.color import RGBColor
import os

SECRET_KEY = os.getenv("ADO_SECRET")
ALIAS = os.getenv("USER_ALIAS")
#This references the PPTX master format for the populated deck.
dir = os.path.dirname(__file__)
PPTXTemplate = os.path.join(dir, 'RedZoneInput.pptx')
exlaimIcon = os.path.join(dir, "exclaim.png")
refreshIcon = os.path.join(dir, "refresh.png")
resourceIcon = os.path.join(dir, "resource.png")

#Functions definitions This is an addition to the code

def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report by providing an outfile')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file')
    return parser.parse_args()

def RZ_Selector():
    response = int(input('Which RedZone would you like to create a Slide deck for(Input number)?\n (1) M365\n (2) Azure\n (3) EMM\n (4) Identity\n (5) Power Platform\n (6) PowerBi\n (7) SDP\n (8) MIP\n (9) D365\n (10) COMM\n:'))

    if response == 1:
        RZTag = 'RZ-M365'
        RZName= 'M365'
        return RZTag, RZName
    elif response == 2:
        RZTag = 'RZ-Azure'
        RZName = 'Azure'
        return RZTag, RZName
    elif response == 3:
        RZTag = 'RZ-EMM'
        RZName = 'EMM'
        return RZTag, RZName
    elif response == 4:
        RZTag = 'RZ-Identity'
        RZName = 'Identity'
        return RZTag, RZName
    elif response == 5:
        RZTag = 'RZ-Power'
        RZName = 'Power Platform'
        return RZTag, RZName
    elif response == 6:
        RZTag = 'RZ-PBI'
        RZName = 'Power Bi'
        return RZTag, RZName
    elif response == 7:
        RZTag = 'RZ-SDP'
        RZName = 'Securing the Developer Pipleine'
        return RZTag, RZName
    elif response == 8:
        RZTag = 'RZ-MIP'
        RZName = 'Microsoft Information Protection'
        return RZTag, RZName
    elif response == 9:
        RZTag = 'RZ-D365'
        RZName = 'Dynamics 365'
        return RZTag, RZName
    elif response == 10:
        RZTag = 'RZ-Comm'
        RZName = 'Commerce'
        return RZTag, RZName
    else:
        print('Invalid entry Try again')
        RZ_Selector()  

  

def API_Pull(RZTag):
    #Wiql Query that is used to post to the REST API and return the Workitems
    wiql_API_url = 'https://dev.azure.com/MicrosoftIT/OneITVSO/_apis/wit/wiql?api-version=5.1'
    
    query = {"query": """SELECT
        [System.Id]
    FROM workitems
    WHERE
        [System.TeamProject] = @project
        AND (
            (
                [System.WorkItemType] = 'Issue'
                AND [System.State] = 'New'
                OR [System.State] = 'Active'
            )
            AND [System.Tags] CONTAINS '{}'
            AND [System.WorkItemType] = 'Issue'
        )
    ORDER BY [System.Id]""".format(RZTag)}

    r = requests.post(wiql_API_url, json = query , auth= HTTPBasicAuth(ALIAS,SECRET_KEY))
    #This is the raw full format JSON from the API Call
    data = r.json()
    #This is filtered down to the Workitems list of Dictionaries ID:URL
    workitems = data['workItems']
    #This uses the defaultdict module to separate the list of dictionaries into 2 lists
    res = defaultdict(list)
    {res[key].append(sub[key]) for sub in workitems for key in sub}
    #Filters the lists to just the list of Workitem Ids. 
    ids = res['id']

    workitem_API_url = 'https://dev.azure.com/MicrosoftIT/OneITVSO/_apis/wit/workitemsbatch?api-version=5.1'


    data = {
      "ids": ids,
      "fields": [
        'System.Id',
        'System.WorkItemType',
        'System.Title',
        'System.AssignedTo',
        'System.State',
        'System.Tags',
        'Microsoft.VSTS.Scheduling.DueDate',
        'System.Description'
      ]
    }

    r = requests.post(workitem_API_url, json = data , auth= HTTPBasicAuth(ALIAS,SECRET_KEY))
    RZone = r.json()

    RZone = RZone['value']
    

    df = pd.json_normalize(RZone)

    #This extracts the status from the Tags fied
    df.loc[df['fields.System.Tags'].str.contains('RZ-Red'),'Status'] = 'RZ-Red'
    df.loc[df['fields.System.Tags'].str.contains('RZ-Yellow'),'Status'] = 'RZ-Yellow'
    df.loc[df['fields.System.Tags'].str.contains('RZ-Green'),'Status'] = 'RZ-Green'
    df.loc[df['fields.System.Tags'].str.contains('RZ-Blue'),'Status'] = 'RZ-Blue'

    #Parse out the HTML
    #Grabs the Description column in the API data
    description = df['fields.System.Description']

    #Instantiates a blank DF object to append to in the loop below
    df3 = pd.DataFrame()

    #Parses through each HTML and appends it to a Dataframe
    for rows in description:
      html = pd.read_html(rows)
      html = html[0]
    
      headowner = html.iloc[1,0]
      headado = html.iloc[2,0]
      headcomments = html.iloc[3,0]
      owner = html.iloc[1,1]
      ado = html.iloc[2,1]
      comments = html.iloc[3,1]
    
      fields = pd.DataFrame(
        {headowner : [owner],
         headado : [ado],
         headcomments : [comments]
        }
      )
      df3 =df3.append(fields)
    
    #df3

    PG_Owner = df3['PG Owner:'].to_list()
    PG_ADO = df3['PG ADO(URL):'].to_list()
    Comments = df3['Comments(Status):'].to_list()

    #This appends and drops uneccesarry columns to the Dataframe. 

    #This appends the new columns
    df['PG Owner'] = PG_Owner
    df['PG ADO'] = PG_ADO
    df['Comments'] = Comments

    #This deletes unecessary columns
    df = df.drop(columns=['id','rev','url','fields.System.WorkItemType','fields.System.State','fields.System.AssignedTo.url','fields.System.AssignedTo._links.avatar.href',
    'fields.System.AssignedTo.id','fields.System.AssignedTo.imageUrl','fields.System.AssignedTo.descriptor','fields.System.Description','fields.System.Tags'])

    #This makes the ADO ID into a URL
    ado_id = df['fields.System.Id'].to_list()
    ado_url = []
    for items in ado_id:
        url = 'https://microsoftit.visualstudio.com/OneITVSO/_workitems/edit/{}'.format(items)
        ado_url.append(url)

    df['MSD ADO'] = ado_url

    #This renames all the columns foeasier referencing
    df = df.rename(columns = {'fields.System.Id':'MSD ADO ID','fields.System.AssignedTo.displayName':'MSD Owner','fields.System.AssignedTo.uniqueName':'MSD Owner Alias', 'fields.System.Title':'Issue','fields.Microsoft.VSTS.Scheduling.DueDate':'Req Date','MSD ADO': 'MSD ADO URL' })

    #This reformats the Date and Time Block to Month and Year
    for index_label, row_series in df.iterrows():
        if type(df.at[index_label , 'Req Date']) == str:
           df.at[index_label , 'Req Date'] = dateutil.parser.parse(row_series['Req Date']).strftime('%m/%y')
        else:
           df.at[index_label , 'Req Date'] = row_series['Req Date']
    #This sorts all the work items from Blue to Red Status.
    df['Status'] = pd.Categorical(df['Status'], ['RZ-Blue', 'RZ-Green', 'RZ-Yellow', 'RZ-Red'])

    df = df.sort_values('Status')

    return df



def create_ppt(input, output, data, RZTitle):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[12]
    title.text = "{} Monthly RedZone".format(RZTitle) #Need to fix this to be custom
    subtitle.text = "Date: {:%B %Y}".format(date.today()) 
        
    slidedata = data[['Issue','Req Date','Status','MSD ADO ID','MSD Owner','PG ADO','PG Owner','Comments']]

    #Numerate the amount of RedZone items to set the number of slides you will need. 
    rz_items = len(data)
    slide_number = (-(-rz_items//7))-1
    #Create New Slide using the Table Template
    slide_start = 0
    while slide_start <= slide_number: 
        redzone_item_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(redzone_item_layout)
        shapes = slide.shapes
        shapes.title.text = 'RedZone Items'
        #Set the Table dimensions in the slide
        rows = cols = 8
        left = Inches(0)
        top = Inches(0.82)
        width = Inches(13.33)
        height = Inches(6.69)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        #Set the Column widths
        table.columns[0].width= Inches(4.0)
        table.columns[1].width= Inches(1.02)
        table.columns[2].width= Inches(0.9)
        table.columns[3].width= Inches(0.89)
        table.columns[4].width= Inches(1.31)
        table.columns[5].width= Inches(0.8)
        table.columns[6].width= Inches(1.03)
        table.columns[7].width= Inches(3.38)
        table.rows[0].height = Inches(0.51)
        #Insert Top row values
        issue_title = table.cell(0,0)
        issue_req_date = table.cell(0,1)
        issue_status = table.cell(0,2)
        issue_id = table.cell(0,3)
        issue_owner = table.cell(0,4)
        issue_pg_id = table.cell(0,5)
        issue_pg_owner = table.cell(0,6)
        issue_comments = table.cell(0,7)
        #Adjust the text formatting on the top row
        #Add top row column values
        issue_title.text = "Issue"
        issue_req_date.text = "Req Date"
        issue_status.text = "Status"
        issue_id.text = "MSD ADO"
        issue_owner.text = "MSD Owner"
        issue_pg_id.text = "PG ADO"
        issue_pg_owner.text = "PG Owner"
        issue_comments.text = "Comments"
        #Adjust the Text formatting on the top row
        columns = [i for i in range(8)]
        for cols in columns:
            table.cell(0, cols).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(0, cols).text_frame.paragraphs[0].font.bold = True
            table.cell(0, cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
        slide_start += 1
        #Insert the moveable icons into the slides. 
        refresh = slide.shapes.add_picture(refreshIcon, left=Inches(6.98), top=Inches(0.08),height=Inches(0.17), width = Inches(0.23))
        exclaimation = slide.shapes.add_picture(exlaimIcon, left=Inches(7.07), top=Inches(0.34),height=Inches(0.28), width = Inches(0.05))
        resource = slide.shapes.add_picture(resourceIcon, left=Inches(5.73), top=Inches(0.03),height=Inches(0.27), width = Inches(0.26))
    #Insert values into the table.
    dataSlide = 1
    redzoneRow = 0
    while (redzoneRow < rz_items):
        populatedSlide = prs.slides[dataSlide]
        table = [shape for shape in populatedSlide.shapes if shape.has_table]
        for rows in range(1,8):
            for cols in range(0,8):
                #This adds hyperlinks to PG ADO Links.
                if 'https:' in str(slidedata.iloc[redzoneRow, cols]):
                    cell = table[0].table.cell(rows,cols)
                    cell_link = cell.text_frame.paragraphs[0].add_run()
                    cell_link.text = 'PG Link' 
                    hlink = cell_link.hyperlink
                    hlink.address = str(slidedata.iloc[redzoneRow, cols])
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Truncates the Character limit of the comments section to 200chars
                elif cols == 7:
                    condense = str(slidedata.iloc[redzoneRow, cols])[:220]
                    table[0].table.cell(rows,cols).text = condense 
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(9)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #This will add Hyperlinks to MSD ADO items.
                elif cols == 3:
                    cell = table[0].table.cell(rows,cols)
                    cell_link = cell.text_frame.paragraphs[0].add_run()
                    cell_link.text = str(slidedata.iloc[redzoneRow, cols]) 
                    hlink = cell_link.hyperlink
                    hlink.address = 'https://microsoftit.visualstudio.com/OneITVSO/_workitems/edit/'+ str(slidedata.iloc[redzoneRow, cols])
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'               
                #Logig for the status cell to be formatted properly.
                elif "RZ-" in str(slidedata.iloc[redzoneRow, cols]):
                    if 'Red' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xFF, 0x4C, 0x4C)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Yellow' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xFF, 0xD3, 0x4C)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Green' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xB2, 0xDE, 0x84)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Blue' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0x4C, 0xC8, 0xF4)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    else:
                        continue
                else:
                    table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the MSD Owner Column 
                table[0].table.cell(rows, 4).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 4).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the PGADO Column 
                table[0].table.cell(rows, 5).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 5).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the PG Owner column         
                table[0].table.cell(rows, 6).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 6).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'                    
            redzoneRow += 1
            if redzoneRow >= rz_items:
                break
            else:
                continue
        dataSlide += 1

    prs.save(output)

if __name__ == "__main__":
    args = parse_args()
    RedZone = RZ_Selector()
    RZTag = RedZone[0]
    RZTitle = RedZone[1]
    data = API_Pull(RZTag)
    create_ppt(PPTXTemplate, args.outfile.name, data, RZTitle)
